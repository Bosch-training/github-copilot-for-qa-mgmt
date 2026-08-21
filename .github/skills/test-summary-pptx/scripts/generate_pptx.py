"""Generate a test-execution summary PPTX from a JSON config, for any test suite (not just AEB).

Usage:
    python3 generate_pptx.py --config path/to/config.json [--output path/to/report.pptx]

See .github/skills/test-summary-pptx/assets/example-config.aeb.json for the config schema.
"""
import argparse
import json

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.chart.data import ChartData
from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION

DARK = RGBColor(0x1F, 0x2A, 0x37)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
GRAY_TEXT = RGBColor(0x60, 0x6B, 0x76)

# Maps a cell's text value (case-insensitive) to a highlight color, used for
# any status-like column (Result, Defect Status, Iteration Result, ...).
STATUS_COLORS = {
    "passed": RGBColor(0x2E, 0xA0, 0x4A),
    "pass": RGBColor(0x2E, 0xA0, 0x4A),
    "correct": RGBColor(0x2E, 0xA0, 0x4A),
    "failed": RGBColor(0xC0, 0x2B, 0x2B),
    "fail": RGBColor(0xC0, 0x2B, 0x2B),
    "open": RGBColor(0xC0, 0x2B, 0x2B),
    "not executed": RGBColor(0x9E, 0x9E, 0x9E),
    "blocked": RGBColor(0x9E, 0x9E, 0x9E),
}


def status_color(value):
    return STATUS_COLORS.get(value.strip().lower())


def add_title(slide, text, subtitle=None):
    box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12.3), Inches(0.8))
    tf = box.text_frame
    tf.text = text
    tf.paragraphs[0].font.size = Pt(30)
    tf.paragraphs[0].font.bold = True
    tf.paragraphs[0].font.color.rgb = DARK
    if subtitle:
        p = tf.add_paragraph()
        p.text = subtitle
        p.font.size = Pt(14)
        p.font.color.rgb = GRAY_TEXT


def add_title_slide(prs, blank, cfg):
    s = prs.slides.add_slide(blank)
    box = s.shapes.add_textbox(Inches(1), Inches(2.6), Inches(11.3), Inches(2))
    tf = box.text_frame
    tf.text = cfg["title"]
    tf.paragraphs[0].font.size = Pt(40)
    tf.paragraphs[0].font.bold = True
    tf.paragraphs[0].font.color.rgb = DARK
    p = tf.add_paragraph()
    p.text = cfg.get("subtitle", "")
    p.font.size = Pt(18)
    p.font.color.rgb = GRAY_TEXT
    p2 = tf.add_paragraph()
    p2.text = f"Report date: {cfg['report_date']}"
    p2.font.size = Pt(14)
    p2.font.color.rgb = GRAY_TEXT


def add_table_slide(prs, blank, heading, note, rows_data, colored_columns,
                     footer=None, col_widths=None, table_height=5.6):
    """Render a heading + table slide. colored_columns: set of column indexes to color by STATUS_COLORS."""
    s = prs.slides.add_slide(blank)
    add_title(s, heading, note)

    rows, cols = len(rows_data), len(rows_data[0])
    table_shape = s.shapes.add_table(
        rows, cols, Inches(0.5), Inches(1.4), Inches(12.3), Inches(table_height))
    table = table_shape.table
    if col_widths:
        for idx, w in enumerate(col_widths):
            table.columns[idx].width = Inches(w)

    for r, row in enumerate(rows_data):
        for c, val in enumerate(row):
            cell = table.cell(r, c)
            cell.text = str(val)
            para = cell.text_frame.paragraphs[0]
            para.font.size = Pt(13) if r == 0 else Pt(12)
            para.font.bold = (r == 0)
            if r == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = DARK
                para.font.color.rgb = WHITE
            elif c in colored_columns:
                color = status_color(str(val))
                if color:
                    para.font.bold = True
                    para.font.color.rgb = color

    if footer:
        note_box = s.shapes.add_textbox(
            Inches(0.5), Inches(1.4 + table_height + 0.1), Inches(12.3), Inches(1.5))
        tf = note_box.text_frame
        tf.word_wrap = True
        tf.text = footer
        tf.paragraphs[0].font.size = Pt(14)
        tf.paragraphs[0].font.italic = True
    return s


def add_summary_slide(prs, blank, cfg):
    summary = cfg["summary"]
    s = prs.slides.add_slide(blank)
    add_title(s, "Summary", cfg.get("subtitle", ""))

    chart_data = ChartData()
    chart_data.categories = summary["donut_labels"]
    chart_data.add_series("Result", tuple(summary["donut_values"]))

    x, y, cx, cy = Inches(0.5), Inches(1.5), Inches(6.0), Inches(5.3)
    graphic_frame = s.shapes.add_chart(XL_CHART_TYPE.DOUGHNUT, x, y, cx, cy, chart_data)
    chart = graphic_frame.chart
    chart.has_legend = True
    chart.legend.position = XL_LEGEND_POSITION.BOTTOM
    chart.legend.include_in_layout = False
    plot = chart.plots[0]
    plot.has_data_labels = True
    plot.data_labels.number_format = '0'
    plot.data_labels.number_format_is_linked = False

    points = plot.series[0].points
    for pt, label in zip(points, summary["donut_labels"]):
        color = status_color(label) or DARK
        pt.format.fill.solid()
        pt.format.fill.fore_color.rgb = color

    bullets = s.shapes.add_textbox(Inches(6.8), Inches(1.6), Inches(6.0), Inches(5.2))
    tf = bullets.text_frame
    tf.word_wrap = True
    first = True
    for text, bold in summary["bullets"]:
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.text = "• " + text
        p.font.size = Pt(15)
        p.font.bold = bool(bold)
        p.space_after = Pt(10)


def build_presentation(cfg):
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank = prs.slide_layouts[6]

    add_title_slide(prs, blank, cfg)

    classification = cfg["classification"]
    add_table_slide(
        prs, blank, "Test Result Classification", classification.get("note"),
        [["Requirement", "Test", "Result"]] + classification["rows"],
        colored_columns={2}, col_widths=[1.6, 8.7, 2.0])

    linkage = cfg["failed_defect_linkage"]
    add_table_slide(
        prs, blank, "Failed Tests — Defect Linkage", linkage.get("note"),
        [["Requirement", "Failure Reason", "Linked Defect", "Defect Status"]] + linkage["rows"],
        colored_columns={3}, col_widths=[1.5, 6.3, 2.0, 2.5],
        footer=linkage.get("footer"), table_height=2.6)

    iterations = cfg["passed_iteration_check"]
    add_table_slide(
        prs, blank, "Passed Tests — Iteration Check", iterations.get("note"),
        [["Requirement", "Test", "Iterations Logged", "Iteration Result"]] + iterations["rows"],
        colored_columns={3}, col_widths=[1.5, 4.3, 2.3, 4.2],
        footer=iterations.get("footer"), table_height=2.4)

    add_summary_slide(prs, blank, cfg)
    return prs


def main():
    parser = argparse.ArgumentParser(description=__doc__)
    parser.add_argument("--config", required=True, help="Path to the report config JSON file")
    parser.add_argument("--output", help="Path to write the .pptx (defaults to config's output_path)")
    args = parser.parse_args()

    with open(args.config) as f:
        cfg = json.load(f)

    output_path = args.output or cfg["output_path"]
    prs = build_presentation(cfg)
    prs.save(output_path)
    print(f"wrote {output_path}")


if __name__ == "__main__":
    main()
